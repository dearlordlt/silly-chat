"""Attachment uploads: images (for vision) and documents (for RAG).

Storage discipline, in order:
  1. Don't fill the disk. Images are recompressed; documents are reduced to text chunks +
     embeddings (the original file is purged after ``doc_file_ttl_days``, the tiny chunks
     live on for ``upload_ttl_days`` so the chat keeps its context). Dedup by content hash;
     per-user + global quotas; TTL sweep.
  2. Owner-only. An upload is readable only by the user who created it.
  3. Cheap to wire into a turn. ``resolve_attachments`` returns image bytes + doc chunks.
"""

from __future__ import annotations

import hashlib
import io
import mimetypes
import uuid
from datetime import timedelta
from pathlib import Path

from fastapi import APIRouter, File, HTTPException, UploadFile, status
from fastapi.responses import FileResponse
from PIL import Image, ImageOps
from pypdf import PdfReader
from sqlalchemy import func
from sqlmodel import Session, select

from app.auth.deps import ApprovedUser, SessionDep
from app.config import ROOT, get_settings
from app.db import engine
from app.embeddings import embed_texts, pack
from app.logging_setup import get_logger
from app.models import DocChunk, Upload, _utcnow

log = get_logger("uploads")
router = APIRouter(prefix="/api/uploads", tags=["uploads"])


def _dir() -> Path:
    d = ROOT / "data" / "uploads"
    d.mkdir(parents=True, exist_ok=True)
    return d


# Document extensions we can extract text from.
_DOC_EXTS = (
    ".pdf", ".docx", ".xlsx", ".pptx",
    ".txt", ".md", ".markdown", ".csv", ".log", ".json", ".xml", ".html", ".htm", ".rtf",
)


def _is_doc(content_type: str, name: str) -> bool:
    ct = (content_type or "").lower()
    return ct == "application/pdf" or ct.startswith("text/") or name.lower().endswith(_DOC_EXTS)


# ---- images -------------------------------------------------------------------------

def _process_image(raw: bytes, max_dim: int) -> bytes:
    """Downscale to max_dim on the long side and re-encode as JPEG (flatten alpha)."""
    img = Image.open(io.BytesIO(raw))
    img = ImageOps.exif_transpose(img)
    if max(img.size) > max_dim:
        img.thumbnail((max_dim, max_dim))
    if img.mode in ("RGBA", "LA", "P"):
        rgba = img.convert("RGBA")
        bg = Image.new("RGB", rgba.size, (255, 255, 255))
        bg.paste(rgba, mask=rgba.split()[-1])
        img = bg
    else:
        img = img.convert("RGB")
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=85, optimize=True)
    return buf.getvalue()


# ---- documents ----------------------------------------------------------------------

def _extract_pdf(raw: bytes) -> str:
    return "\n\n".join((p.extract_text() or "") for p in PdfReader(io.BytesIO(raw)).pages)


def _extract_docx(raw: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(raw))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_xlsx(raw: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts)


def _extract_pptx(raw: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(raw))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def _extract_text(raw: bytes, content_type: str, name: str) -> str:
    """Extract plain text from a supported document, capped to doc_max_chars."""
    lower = name.lower()
    try:
        if lower.endswith(".pdf") or (content_type or "").lower() == "application/pdf":
            text = _extract_pdf(raw)
        elif lower.endswith(".docx"):
            text = _extract_docx(raw)
        elif lower.endswith(".xlsx"):
            text = _extract_xlsx(raw)
        elif lower.endswith(".pptx"):
            text = _extract_pptx(raw)
        else:
            text = raw.decode("utf-8", errors="replace")
    except Exception as exc:
        log.warning("extract failed for %s: %s", name, exc)
        text = ""
    cap = get_settings().limits.doc_max_chars
    return text[:cap]


def _chunk(text: str, size: int, overlap: int) -> list[str]:
    text = text.strip()
    if not text:
        return []
    chunks: list[str] = []
    step = max(1, size - overlap)
    for start in range(0, len(text), step):
        piece = text[start : start + size].strip()
        if piece:
            chunks.append(piece)
        if start + size >= len(text):
            break
    return chunks


# ---- quotas / lifecycle -------------------------------------------------------------

def _user_used(session: Session, uid: int) -> int:
    return session.exec(select(func.coalesce(func.sum(Upload.size), 0)).where(Upload.user_id == uid)).one()


def _global_used(session: Session) -> int:
    return session.exec(select(func.coalesce(func.sum(Upload.size), 0))).one()


def _delete_upload(session: Session, up: Upload) -> None:
    for ch in session.exec(select(DocChunk).where(DocChunk.upload_id == up.id)).all():
        session.delete(ch)
    session.delete(up)


def _gc_orphan_files(session: Session) -> None:
    referenced = set(session.exec(select(Upload.sha256)).all())
    for f in _dir().glob("*.*"):
        if f.stem not in referenced:
            f.unlink(missing_ok=True)


def _evict_until_room(session: Session, incoming: int, limit: int) -> None:
    used = _global_used(session)
    if used + incoming <= limit:
        return
    for up in session.exec(select(Upload).order_by(Upload.last_used_at)).all():
        if used + incoming <= limit:
            break
        used -= up.size
        _delete_upload(session, up)
    session.commit()
    _gc_orphan_files(session)


def sweep_expired() -> int:
    """Purge doc originals past their (short) file TTL, then delete fully-expired uploads
    (rows + chunks + files). Called at startup."""
    cfg = get_settings().limits
    now = _utcnow()
    doc_file_cutoff = now - timedelta(days=cfg.doc_file_ttl_days)
    full_cutoff = now - timedelta(days=cfg.upload_ttl_days)
    with Session(engine) as session:
        # 1) drop the heavy original doc files; keep the row + chunks (context survives).
        for up in session.exec(
            select(Upload).where(Upload.kind == "doc", Upload.created_at < doc_file_cutoff)
        ).all():
            (_dir() / f"{up.sha256}.{up.ext}").unlink(missing_ok=True)
        # 2) fully expire old uploads.
        expired = session.exec(select(Upload).where(Upload.created_at < full_cutoff)).all()
        for up in expired:
            _delete_upload(session, up)
        session.commit()
        _gc_orphan_files(session)
    if expired:
        log.info("upload sweep: removed %d expired", len(expired))
    return len(expired)


# ---- endpoints ----------------------------------------------------------------------

@router.post("")
async def create_upload(
    user: ApprovedUser, session: SessionDep, file: UploadFile = File(...)
) -> dict:
    cfg = get_settings().limits
    raw = await file.read()
    name = file.filename or "file"
    is_image = (file.content_type or "").startswith("image/")
    is_doc = _is_doc(file.content_type or "", name)
    if not (is_image or is_doc):
        raise HTTPException(status.HTTP_415_UNSUPPORTED_MEDIA_TYPE, "unsupported file type")

    cap = (cfg.upload_max_mb if is_image else cfg.doc_max_mb) * 1024 * 1024
    if len(raw) > cap:
        raise HTTPException(status.HTTP_413_CONTENT_TOO_LARGE, "file too large")

    if is_image:
        try:
            data = _process_image(raw, cfg.upload_image_max_dim)
        except Exception:
            raise HTTPException(status.HTTP_400_BAD_REQUEST, "could not read that image")
        kind, mime, ext = "image", "image/jpeg", "jpg"
    else:
        data = raw
        kind = "doc"
        mime = mimetypes.guess_type(name)[0] or "application/octet-stream"
        ext = name.rsplit(".", 1)[-1].lower() if "." in name else "txt"

    if _user_used(session, user.id) + len(data) > cfg.upload_user_quota_mb * 1024 * 1024:
        raise HTTPException(status.HTTP_413_CONTENT_TOO_LARGE, "you've reached your upload quota")
    _evict_until_room(session, len(data), cfg.upload_global_quota_mb * 1024 * 1024)

    sha = hashlib.sha256(data).hexdigest()
    path = _dir() / f"{sha}.{ext}"
    if not path.exists():
        path.write_bytes(data)
    up = Upload(
        id=uuid.uuid4().hex, user_id=user.id, sha256=sha, kind=kind,
        mime=mime, ext=ext, size=len(data), name=name,
    )
    session.add(up)
    session.commit()

    if kind == "doc":
        text = _extract_text(raw, file.content_type or "", name)
        chunks = _chunk(text, cfg.doc_chunk_chars, cfg.doc_chunk_overlap)
        if not chunks:
            raise HTTPException(status.HTTP_400_BAD_REQUEST, "no readable text in that document")
        vecs = await embed_texts(chunks)
        for i, (ctext, vec) in enumerate(zip(chunks, vecs)):
            session.add(DocChunk(upload_id=up.id, idx=i, text=ctext, embedding=pack(vec)))
        session.commit()
        log.info("doc %s: %d chunks embedded", name, len(chunks))
        return {"id": up.id, "kind": kind, "name": name, "mime": mime, "chunks": len(chunks)}

    return {"id": up.id, "kind": kind, "name": name, "mime": mime}


@router.get("/{uid}")
def serve_upload(uid: str, user: ApprovedUser, session: SessionDep) -> FileResponse:
    up = session.get(Upload, uid)
    if up is None or up.user_id != user.id:
        raise HTTPException(status.HTTP_404_NOT_FOUND, "not found")
    path = _dir() / f"{up.sha256}.{up.ext}"
    if not path.exists():
        raise HTTPException(status.HTTP_404_NOT_FOUND, "expired")
    up.last_used_at = _utcnow()
    session.add(up)
    session.commit()
    return FileResponse(path, media_type=up.mime, headers={"Cache-Control": "private, max-age=86400"})


def resolve_attachments(
    session: Session, ids: list[str], user_id: int, include_docs: bool = True
) -> tuple[list[tuple[str, bytes]], list[tuple[str, bytes]]]:
    """Resolve upload ids (owner-checked) into (images, doc_chunks).

    images: list of (mime, bytes) for the vision tool.
    doc_chunks: list of (text, embedding_bytes) for RAG (empty unless include_docs).
    """
    images: list[tuple[str, bytes]] = []
    doc_chunks: list[tuple[str, bytes]] = []
    for aid in ids:
        up = session.get(Upload, aid)
        if up is None or up.user_id != user_id:
            continue
        if up.kind == "image":
            path = _dir() / f"{up.sha256}.{up.ext}"
            if path.exists():
                up.last_used_at = _utcnow()
                session.add(up)
                images.append((up.mime, path.read_bytes()))
        elif up.kind == "doc" and include_docs:
            up.last_used_at = _utcnow()
            session.add(up)
            for ch in session.exec(select(DocChunk).where(DocChunk.upload_id == up.id)).all():
                doc_chunks.append((ch.text, ch.embedding))
    session.commit()
    return images, doc_chunks
